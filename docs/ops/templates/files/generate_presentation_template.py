#!/usr/bin/env python3
"""
Generate AMP Presentation Template (PowerPoint)

Creates a professional presentation template for real estate project marketing plans
with pre-structured slides for: Overview, Target Audience, Message, Channels, Budget,
Creative Direction, KPIs, and Tracking.

Requirements:
    pip install python-pptx

Usage:
    python generate_presentation_template.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation_template():
    """Create the presentation template with all required slides."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    BRAND_PRIMARY = RGBColor(41, 128, 185)  # Blue
    BRAND_SECONDARY = RGBColor(52, 73, 94)  # Dark Blue
    ACCENT_COLOR = RGBColor(231, 76, 60)    # Red
    TEXT_DARK = RGBColor(44, 62, 80)        # Dark Gray
    TEXT_LIGHT = RGBColor(236, 240, 241)    # Light Gray
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_PRIMARY
    
    # Add title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "‡πÅ‡∏ú‡∏ô‡∏Å‡∏≤‡∏£‡∏ï‡∏•‡∏≤‡∏î‡∏î‡∏¥‡∏à‡∏¥‡∏ó‡∏±‡∏•\nDigital Marketing Plan"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_LIGHT
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "[‡∏ä‡∏∑‡πà‡∏≠‡πÇ‡∏Ñ‡∏£‡∏á‡∏Å‡∏≤‡∏£ / Project Name]"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = TEXT_LIGHT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add presenter info
    info_box = slide1.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    info_frame = info_box.text_frame
    info_frame.text = "‡∏ô‡∏≥‡πÄ‡∏™‡∏ô‡∏≠‡πÇ‡∏î‡∏¢: [‡∏ä‡∏∑‡πà‡∏≠‡∏ó‡∏µ‡∏°/‡∏ö‡∏£‡∏¥‡∏©‡∏±‡∏ó] | ‡∏ß‡∏±‡∏ô‡∏ó‡∏µ‡πà: [DD/MM/YYYY]"
    info_para = info_frame.paragraphs[0]
    info_para.font.size = Pt(18)
    info_para.font.color.rgb = TEXT_LIGHT
    info_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Overview / Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide2, "1. Overview / ‡∏†‡∏≤‡∏û‡∏£‡∏ß‡∏°‡πÇ‡∏Ñ‡∏£‡∏á‡∏Å‡∏≤‡∏£", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üìã Campaign Overview", 24, TEXT_DARK, True)
    add_bullet_point(tf, "‚Ä¢ Objective: [‡πÄ‡∏õ‡πâ‡∏≤‡∏´‡∏°‡∏≤‡∏¢‡∏´‡∏•‡∏±‡∏Å‡∏Ç‡∏≠‡∏á‡πÅ‡∏Ñ‡∏°‡πÄ‡∏õ‡∏ç]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Duration: [‡∏£‡∏∞‡∏¢‡∏∞‡πÄ‡∏ß‡∏•‡∏≤]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Total Budget: ‡∏øXXX,XXX/‡πÄ‡∏î‡∏∑‡∏≠‡∏ô", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "üè† ‡πÇ‡∏Ñ‡∏£‡∏á‡∏Å‡∏≤‡∏£ / Project Details", 24, TEXT_DARK, True)
    add_bullet_point(tf, "‚Ä¢ ‡∏ä‡∏∑‡πà‡∏≠‡πÇ‡∏Ñ‡∏£‡∏á‡∏Å‡∏≤‡∏£: [Project Name]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ ‡∏õ‡∏£‡∏∞‡πÄ‡∏†‡∏ó: [Condo/Villa/House]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ ‡∏£‡∏≤‡∏Ñ‡∏≤: ‡∏øX.X - ‡∏øXX ‡∏•‡πâ‡∏≤‡∏ô", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ ‡∏ó‡∏≥‡πÄ‡∏•: [Location, Pattaya]", 18, TEXT_DARK, False)
    
    # Slide 3: Target Audience
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide3, "2. Target Audience / ‡∏Å‡∏•‡∏∏‡πà‡∏°‡πÄ‡∏õ‡πâ‡∏≤‡∏´‡∏°‡∏≤‡∏¢", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üéØ Primary Audience / ‡∏Å‡∏•‡∏∏‡πà‡∏°‡πÄ‡∏õ‡πâ‡∏≤‡∏´‡∏°‡∏≤‡∏¢‡∏´‡∏•‡∏±‡∏Å", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "Persona 1: [‡∏ä‡∏∑‡πà‡∏≠ Persona]", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Age: XX-XX years", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Income: ‡∏øXXX,XXX+/month", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Occupation: [‡∏≠‡∏≤‡∏ä‡∏µ‡∏û]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Motivation: [‡πÅ‡∏£‡∏á‡∏à‡∏π‡∏á‡πÉ‡∏à]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Pain Points: [‡∏õ‡∏±‡∏ç‡∏´‡∏≤‡∏ó‡∏µ‡πà‡∏ï‡πâ‡∏≠‡∏á‡πÅ‡∏Å‡πâ]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "üìä Behavioral Insights / ‡∏û‡∏§‡∏ï‡∏¥‡∏Å‡∏£‡∏£‡∏°", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ [Insight 1]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ [Insight 2]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ [Insight 3]", 18, TEXT_DARK, False)
    
    # Slide 4: Message & Creative Direction
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide4, "3. Message & Creative Direction / ‡πÅ‡∏ô‡∏ß‡∏Ñ‡∏¥‡∏î‡πÅ‡∏•‡∏∞‡∏™‡∏∑‡πà‡∏≠‡∏™‡∏≤‡∏£", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üí¨ Main Message / ‡∏Ç‡πâ‡∏≠‡∏Ñ‡∏ß‡∏≤‡∏°‡∏´‡∏•‡∏±‡∏Å", 24, TEXT_DARK, True)
    add_bullet_point(tf, '‚Ä¢ Thai: "[‡∏™‡πÇ‡∏•‡πÅ‡∏Å‡∏ô‡∏†‡∏≤‡∏©‡∏≤‡πÑ‡∏ó‡∏¢]"', 18, TEXT_DARK, False)
    add_bullet_point(tf, '‚Ä¢ English: "[English Slogan]"', 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "üé® Creative Direction / ‡∏ó‡∏¥‡∏®‡∏ó‡∏≤‡∏á‡∏Ñ‡∏£‡∏µ‡πÄ‡∏≠‡∏ó‡∏µ‡∏ü", 24, TEXT_DARK, True)
    add_bullet_point(tf, "‚Ä¢ Visual Style: [e.g., Luxury & Premium, Modern & Clean]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Tone of Voice: [e.g., Professional, Aspirational]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Key Visual Elements: [‡∏™‡∏µ, ‡∏ü‡∏≠‡∏ô‡∏ï‡πå, ‡∏™‡πÑ‡∏ï‡∏•‡πå‡∏†‡∏≤‡∏û]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "üéØ Key Messages by Audience / ‡∏Ç‡πâ‡∏≠‡∏Ñ‡∏ß‡∏≤‡∏°‡∏ï‡∏≤‡∏°‡∏Å‡∏•‡∏∏‡πà‡∏°‡πÄ‡∏õ‡πâ‡∏≤‡∏´‡∏°‡∏≤‡∏¢", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Audience 1: [Message]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Audience 2: [Message]", 18, TEXT_DARK, False)
    
    # Slide 5: Channels Strategy
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide5, "4. Channels / ‡∏ä‡πà‡∏≠‡∏á‡∏ó‡∏≤‡∏á", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üì¢ Channel Mix / ‡∏ä‡πà‡∏≠‡∏á‡∏ó‡∏≤‡∏á‡πÇ‡∏Ü‡∏©‡∏ì‡∏≤", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "1. Google Ads (35% ‡∏á‡∏ö)", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Search: Branded + Generic keywords", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Display: Retargeting website visitors", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Target: [X] leads, CPL < ‡∏øXXX", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "2. Facebook / Instagram (35% ‡∏á‡∏ö)", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Lead Generation: Instant Forms", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Traffic: Drive to landing page", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Audience: Lookalike + Interests + Retargeting", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "3. LINE Ads (15% ‡∏á‡∏ö)", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Talk Head View + Display Ads", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Focus: Retargeting warm audiences", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "4. TikTok (15% ‡∏á‡∏ö)", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ In-Feed Ads: Property tours & lifestyle", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Objective: Brand awareness & engagement", 18, TEXT_DARK, False)
    
    # Slide 6: Budget
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide6, "5. Budget / ‡∏á‡∏ö‡∏õ‡∏£‡∏∞‡∏°‡∏≤‡∏ì", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide6.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üí∞ Monthly Budget Allocation / ‡∏Å‡∏≤‡∏£‡∏à‡∏±‡∏î‡∏™‡∏£‡∏£‡∏á‡∏ö‡∏£‡∏≤‡∏¢‡πÄ‡∏î‡∏∑‡∏≠‡∏ô", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "Advertising Spend:", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Google Ads: ‡∏øXX,XXX (35%)", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Facebook/Instagram: ‡∏øXX,XXX (35%)", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ LINE Ads: ‡∏øX,XXX (15%)", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ TikTok: ‡∏øX,XXX (15%)", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Advertising Subtotal: ‡∏øXX,XXX", 18, ACCENT_COLOR, True)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "Other Costs:", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Content Production: ‡∏øX,XXX", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Tools & Software: ‡∏øX,XXX", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "üìä Grand Total: ‡∏øXX,XXX/month", 22, ACCENT_COLOR, True)
    
    # Slide 7: KPIs
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide7, "6. KPIs / ‡∏ï‡∏±‡∏ß‡∏ä‡∏µ‡πâ‡∏ß‡∏±‡∏î‡∏Ñ‡∏ß‡∏≤‡∏°‡∏™‡∏≥‡πÄ‡∏£‡πá‡∏à", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide7.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üìä Primary KPIs / ‡∏ï‡∏±‡∏ß‡∏ä‡∏µ‡πâ‡∏ß‡∏±‡∏î‡∏´‡∏•‡∏±‡∏Å", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "‚Ä¢ Leads per Month: [X]+ qualified leads", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Cost Per Lead (CPL): < ‡∏øXXX", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Lead Quality Score: [X]/10 average", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Lead-to-Appointment Rate: > XX%", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Return on Ad Spend (ROAS): > X:1", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "üìà Secondary KPIs / ‡∏ï‡∏±‡∏ß‡∏ä‡∏µ‡πâ‡∏ß‡∏±‡∏î‡∏£‡∏≠‡∏á", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "‚Ä¢ Click-Through Rate (CTR): > 1.5%", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Landing Page Conversion Rate: > 5%", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Cost Per Click (CPC): < ‡∏ø20", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Impression Share: > 60%", 18, TEXT_DARK, False)
    
    # Slide 8: Tracking & Reporting
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide8, "7. Tracking / ‡∏Å‡∏≤‡∏£‡∏ï‡∏¥‡∏î‡∏ï‡∏≤‡∏°‡∏ú‡∏•", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide8.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üìä Tracking Setup / ‡∏£‡∏∞‡∏ö‡∏ö‡∏ï‡∏¥‡∏î‡∏ï‡∏≤‡∏°‡∏ú‡∏•", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "Tools & Platforms:", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Google Analytics 4: Website traffic & conversions", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Google Tag Manager: Event tracking", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Meta Pixel: Facebook/Instagram tracking", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Call tracking: Phone lead attribution", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ CRM Integration: Lead management", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "üìÖ Reporting Schedule / ‡∏Å‡∏≥‡∏´‡∏ô‡∏î‡∏Å‡∏≤‡∏£‡∏£‡∏≤‡∏¢‡∏á‡∏≤‡∏ô", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Weekly: Quick performance dashboard", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Monthly: Full report with insights & recommendations", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Quarterly: Strategic review & planning", 18, TEXT_DARK, False)
    
    # Slide 9: Timeline
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide9, "8. Timeline / ‡∏Å‡∏≥‡∏´‡∏ô‡∏î‡∏Å‡∏≤‡∏£", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide9.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "üìÖ Campaign Timeline / ‡πÅ‡∏ú‡∏ô‡∏á‡∏≤‡∏ô", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "Month 1: Setup & Launch", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Week 1-2: Campaign setup, creative production", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Week 3: Soft launch & testing", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Week 4: Full launch", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "Month 2: Testing & Optimization", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ A/B testing audiences & creatives", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Strategy adjustment based on data", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "Month 3: Scaling", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Scale winning campaigns", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Expand to additional audiences", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "Month 4+: Mature & Optimize", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "‚Ä¢ Continue optimization & creative refresh", 18, TEXT_DARK, False)
    add_bullet_point(tf, "‚Ä¢ Consistent ROI & predictable results", 18, TEXT_DARK, False)
    
    # Slide 10: Next Steps / Q&A
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide10, "Next Steps / ‡∏Ç‡∏±‡πâ‡∏ô‡∏ï‡∏≠‡∏ô‡∏ñ‡∏±‡∏î‡πÑ‡∏õ", BRAND_PRIMARY, TEXT_LIGHT)
    
    content_box = slide10.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    add_bullet_point(tf, "‚úÖ Immediate Actions / ‡∏î‡∏≥‡πÄ‡∏ô‡∏¥‡∏ô‡∏Å‡∏≤‡∏£‡∏ó‡∏±‡∏ô‡∏ó‡∏µ", 24, TEXT_DARK, True)
    add_bullet_point(tf, "", 12, TEXT_DARK, False)
    
    add_bullet_point(tf, "1. Approve Strategy & Budget", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "   ‚Ä¢ Confirm budget allocation", 18, TEXT_DARK, False)
    add_bullet_point(tf, "   ‚Ä¢ Sign off on creative direction", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "2. Kickoff Meeting", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "   ‚Ä¢ Date: [Proposed date]", 18, TEXT_DARK, False)
    add_bullet_point(tf, "   ‚Ä¢ Finalize details & assign responsibilities", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "3. Campaign Setup (Week 1-2)", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "   ‚Ä¢ Landing page development", 18, TEXT_DARK, False)
    add_bullet_point(tf, "   ‚Ä¢ Ad account setup & creative production", 18, TEXT_DARK, False)
    add_bullet_point(tf, "", 10, TEXT_DARK, False)
    
    add_bullet_point(tf, "4. Launch (Week 3)", 20, BRAND_SECONDARY, True)
    add_bullet_point(tf, "   ‚Ä¢ Soft launch ‚Üí Monitor ‚Üí Full launch", 18, TEXT_DARK, False)
    
    # Add contact info at bottom
    contact_box = slide10.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
    contact_tf = contact_box.text_frame
    contact_tf.text = "üìû Contact: [Your Name] | Email: [email] | LINE: [@lineid]"
    contact_para = contact_tf.paragraphs[0]
    contact_para.font.size = Pt(14)
    contact_para.font.color.rgb = TEXT_DARK
    contact_para.alignment = PP_ALIGN.CENTER
    
    return prs


def add_slide_header(slide, text, bg_color, text_color):
    """Add a header bar to a slide."""
    # Add header background
    header_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = bg_color
    header_bg.line.color.rgb = bg_color
    
    # Add header text
    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = header_text.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = text_color


def add_bullet_point(text_frame, text, font_size, color, bold=False):
    """Add a bullet point to a text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = 0
    return p


def main():
    """Main function to generate the presentation."""
    print("üé® Generating AMP Presentation Template...")
    
    # Create presentation
    prs = create_presentation_template()
    
    # Save presentation
    output_file = "AMP_Presentation_Template.pptx"
    prs.save(output_file)
    
    print(f"‚úÖ Presentation template created: {output_file}")
    print(f"üìä Total slides: {len(prs.slides)}")
    print("\nüìã Slides included:")
    print("   1. Title Slide")
    print("   2. Overview / Executive Summary")
    print("   3. Target Audience")
    print("   4. Message & Creative Direction")
    print("   5. Channels Strategy")
    print("   6. Budget")
    print("   7. KPIs")
    print("   8. Tracking & Reporting")
    print("   9. Timeline")
    print("  10. Next Steps / Q&A")
    print("\nüí° Tip: Open with PowerPoint, Google Slides, or LibreOffice Impress")


if __name__ == "__main__":
    main()
